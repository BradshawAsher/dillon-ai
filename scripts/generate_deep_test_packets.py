import os
import json
import subprocess
import tempfile
from datetime import datetime, timezone
from email.message import EmailMessage

import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether

ROOT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEALS_DIR = os.path.join(ROOT_DIR, "test_sets", "deals")
GT_DIR = os.path.join(ROOT_DIR, "test_sets", "ground_truth")
RESULTS_DIR = os.path.join(ROOT_DIR, "test_sets", "results")

os.makedirs(DEALS_DIR, exist_ok=True)
os.makedirs(GT_DIR, exist_ok=True)
os.makedirs(RESULTS_DIR, exist_ok=True)

# -------------------------------------------------------------------------
# STYLING HELPERS FOR EXCEL
# -------------------------------------------------------------------------
HEADER_NAVY = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
HEADER_FONT = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
SUBHEADER_FILL = PatternFill(start_color="DBEAFE", end_color="DBEAFE", fill_type="solid")
SUBHEADER_FONT = Font(name="Calibri", size=11, bold=True, color="1E3A8A")
TOTAL_FILL = PatternFill(start_color="F1F5F9", end_color="F1F5F9", fill_type="solid")
TOTAL_FONT = Font(name="Calibri", size=11, bold=True, color="0F172A")
REGULAR_FONT = Font(name="Calibri", size=10, color="334155")
ALERT_RED_FILL = PatternFill(start_color="FEE2E2", end_color="FEE2E2", fill_type="solid")
ALERT_RED_FONT = Font(name="Calibri", size=10, bold=True, color="991B1B")
ALERT_YELLOW_FILL = PatternFill(start_color="FEF3C7", end_color="FEF3C7", fill_type="solid")
ALERT_YELLOW_FONT = Font(name="Calibri", size=10, bold=True, color="92400E")

BORDER_THIN = Border(
    left=Side(style='thin', color='CBD5E1'),
    right=Side(style='thin', color='CBD5E1'),
    top=Side(style='thin', color='CBD5E1'),
    bottom=Side(style='thin', color='CBD5E1')
)
BORDER_DOUBLE_BOTTOM = Border(
    left=Side(style='thin', color='CBD5E1'),
    right=Side(style='thin', color='CBD5E1'),
    top=Side(style='thin', color='CBD5E1'),
    bottom=Side(style='double', color='0F172A')
)

def format_excel_tab(ws, title, subtitle, headers, rows, total_indices=None, alert_indices=None, is_currency_cols=None):
    if total_indices is None:
        total_indices = []
    if alert_indices is None:
        alert_indices = []
    if is_currency_cols is None:
        is_currency_cols = range(1, len(headers))
        
    ws.views.sheetView[0].showGridLines = True
    
    # Title Block
    ws.merge_cells("A1:H1")
    t_cell = ws["A1"]
    t_cell.value = title
    t_cell.font = Font(name="Calibri", size=13, bold=True, color="1E3A8A")
    t_cell.alignment = Alignment(vertical="center")
    ws.row_dimensions[1].height = 24
    
    ws.merge_cells("A2:H2")
    s_cell = ws["A2"]
    s_cell.value = subtitle
    s_cell.font = Font(name="Calibri", size=10, italic=True, color="64748B")
    s_cell.alignment = Alignment(vertical="center")
    ws.row_dimensions[2].height = 18
    
    # Headers
    ws.row_dimensions[4].height = 22
    for c_idx, h in enumerate(headers, 1):
        cell = ws.cell(row=4, column=c_idx, value=h)
        cell.fill = HEADER_NAVY
        cell.font = HEADER_FONT
        cell.alignment = Alignment(horizontal="center" if c_idx > 1 else "left", vertical="center")
        cell.border = BORDER_THIN
        
    # Rows
    current_row = 5
    for r_idx, row_data in enumerate(rows):
        ws.row_dimensions[current_row].height = 20
        is_total = r_idx in total_indices
        is_alert = r_idx in alert_indices
        
        for c_idx, val in enumerate(row_data, 1):
            cell = ws.cell(row=current_row, column=c_idx, value=val)
            cell.border = BORDER_DOUBLE_BOTTOM if is_total else BORDER_THIN
            
            if is_total:
                cell.fill = TOTAL_FILL
                cell.font = TOTAL_FONT
            elif is_alert:
                cell.fill = ALERT_RED_FILL
                cell.font = ALERT_RED_FONT
            else:
                cell.font = REGULAR_FONT
                
            if c_idx == 1:
                cell.alignment = Alignment(horizontal="left", vertical="center")
            else:
                cell.alignment = Alignment(horizontal="right", vertical="center")
                if isinstance(val, (int, float)) and c_idx in is_currency_cols:
                    if isinstance(val, float) and abs(val) < 1.0:
                        cell.number_format = '0.0%'
                    else:
                        cell.number_format = '$#,##0'
        current_row += 1
        
    # Auto column widths
    for col in ws.columns:
        col_letter = get_column_letter(col[0].column)
        max_len = max(len(str(cell.value or '')) for cell in col)
        ws.column_dimensions[col_letter].width = max(max_len + 4, 15)


# -------------------------------------------------------------------------
# PDF GENERATION HELPER (REPORTLAB - MULTI-PAGE STYLED)
# -------------------------------------------------------------------------
def generate_branded_pdf(filepath, title, subtitle, metadata_dict, sections):
    doc = SimpleDocTemplate(
        filepath,
        pagesize=letter,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontSize=18,
        leading=22,
        textColor=colors.HexColor("#1E3A8A"),
        spaceAfter=4
    )
    subtitle_style = ParagraphStyle(
        'DocSubtitle',
        parent=styles['Normal'],
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#64748B"),
        spaceAfter=12
    )
    h2_style = ParagraphStyle(
        'SectionH2',
        parent=styles['Heading2'],
        fontSize=12,
        leading=15,
        textColor=colors.HexColor("#1E3A8A"),
        spaceBefore=14,
        spaceAfter=6
    )
    h3_style = ParagraphStyle(
        'SectionH3',
        parent=styles['Heading3'],
        fontSize=10.5,
        leading=13,
        textColor=colors.HexColor("#0F172A"),
        spaceBefore=8,
        spaceAfter=4
    )
    body_style = ParagraphStyle(
        'BodyDark',
        parent=styles['Normal'],
        fontSize=9,
        leading=13,
        textColor=colors.HexColor("#1E293B"),
        spaceAfter=6
    )
    callout_style = ParagraphStyle(
        'CalloutText',
        parent=styles['Normal'],
        fontSize=8.5,
        leading=12,
        textColor=colors.HexColor("#0F172A")
    )
    
    story = []
    story.append(Paragraph(title, title_style))
    story.append(Paragraph(subtitle, subtitle_style))
    
    # Metadata Header Box
    meta_data = []
    keys = list(metadata_dict.keys())
    for i in range(0, len(keys), 2):
        k1 = keys[i]
        v1 = metadata_dict[k1]
        k2 = keys[i+1] if i+1 < len(keys) else ""
        v2 = metadata_dict[k2] if i+1 < len(keys) else ""
        meta_data.append([
            Paragraph(f"<b>{k1}:</b>", body_style),
            Paragraph(str(v1), body_style),
            Paragraph(f"<b>{k2}:</b>" if k2 else "", body_style),
            Paragraph(str(v2) if v2 else "", body_style),
        ])
    
    t_meta = Table(meta_data, colWidths=[110, 155, 110, 165])
    t_meta.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor("#F8FAFC")),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor("#CBD5E1")),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 4),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
    ]))
    story.append(t_meta)
    story.append(Spacer(1, 10))
    
    # Sections
    for sec in sections:
        if sec.get('page_break_before', False):
            story.append(PageBreak())
            
        story.append(Paragraph(sec['heading'], h2_style))
        for p in sec.get('paragraphs', []):
            story.append(Paragraph(p, body_style))
            
        if 'subsections' in sec:
            for sub in sec['subsections']:
                story.append(Paragraph(sub['title'], h3_style))
                for sp in sub.get('paragraphs', []):
                    story.append(Paragraph(sp, body_style))
            
        if 'table' in sec:
            t_info = sec['table']
            tbl_data = []
            tbl_data.append([Paragraph(f"<b>{h}</b>", ParagraphStyle('TH', parent=body_style, textColor=colors.white, alignment=1 if idx > 0 else 0)) for idx, h in enumerate(t_info['headers'])])
            for row in t_info['rows']:
                row_cells = []
                for c_idx, cell_val in enumerate(row):
                    align = 2 if c_idx > 0 else 0
                    row_cells.append(Paragraph(str(cell_val), ParagraphStyle('TD', parent=body_style, alignment=align)))
                tbl_data.append(row_cells)
            
            t_obj = Table(tbl_data, colWidths=t_info.get('colWidths', None))
            t_obj.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#1E3A8A")),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#E2E8F0")),
                ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor("#F8FAFC")]),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('TOPPADDING', (0,0), (-1,-1), 4),
                ('BOTTOMPADDING', (0,0), (-1,-1), 4),
            ]))
            story.append(t_obj)
            story.append(Spacer(1, 8))
            
        if 'callout' in sec:
            c_text = sec['callout']['text']
            c_type = sec['callout'].get('type', 'info')
            bg_color = colors.HexColor("#EFF6FF") if c_type == 'info' else (colors.HexColor("#FEF2F2") if c_type == 'danger' else colors.HexColor("#FFFBEB"))
            border_color = colors.HexColor("#3B82F6") if c_type == 'info' else (colors.HexColor("#EF4444") if c_type == 'danger' else colors.HexColor("#F59E0B"))
            
            c_tbl = Table([[Paragraph(c_text, callout_style)]], colWidths=[540])
            c_tbl.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,-1), bg_color),
                ('BOX', (0,0), (-1,-1), 1, border_color),
                ('LEFTPADDING', (0,0), (-1,-1), 8),
                ('RIGHTPADDING', (0,0), (-1,-1), 8),
                ('TOPPADDING', (0,0), (-1,-1), 6),
                ('BOTTOMPADDING', (0,0), (-1,-1), 6),
            ]))
            story.append(c_tbl)
            story.append(Spacer(1, 8))
            
    doc.build(story)


# -------------------------------------------------------------------------
# PPTX GENERATION HELPER (16:9 MULTI-SLIDE)
# -------------------------------------------------------------------------
def generate_branded_pptx(filepath, title, subtitle, slides_data):
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(5.625)
    blank_layout = prs.slide_layouts[6]
    
    # Title Slide
    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PptxRGBColor(30, 58, 138)
    bg.line.fill.background()
    
    tx_box = title_slide.shapes.add_textbox(PptxInches(1), PptxInches(1.8), PptxInches(8), PptxInches(2))
    tf = tx_box.text_frame
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = PptxPt(30)
    p1.font.bold = True
    p1.font.color.rgb = PptxRGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = PptxPt(15)
    p2.font.color.rgb = PptxRGBColor(191, 219, 254)
    
    # Content Slides
    for s_idx, s in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank_layout)
        
        banner = slide.shapes.add_shape(1, 0, 0, prs.slide_width, PptxInches(0.85))
        banner.fill.solid()
        banner.fill.fore_color.rgb = PptxRGBColor(15, 23, 42)
        banner.line.fill.background()
        
        h_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(0.12), PptxInches(8.8), PptxInches(0.6))
        tf_h = h_box.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = s['title']
        p_h.font.size = PptxPt(19)
        p_h.font.bold = True
        p_h.font.color.rgb = PptxRGBColor(255, 255, 255)
        
        body_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(1.05), PptxInches(8.8), PptxInches(4.1))
        tf_b = body_box.text_frame
        tf_b.word_wrap = True
        
        for b_idx, bullet in enumerate(s.get('bullets', [])):
            p = tf_b.add_paragraph() if b_idx > 0 else tf_b.paragraphs[0]
            p.text = f"•  {bullet}"
            p.font.size = PptxPt(12.5)
            p.font.color.rgb = PptxRGBColor(51, 65, 85)
            p.space_after = PptxPt(8)
            
    prs.save(filepath)


# -------------------------------------------------------------------------
# DOCX GENERATION HELPER (MULTI-PAGE COMPREHENSIVE)
# -------------------------------------------------------------------------
def generate_branded_docx(filepath, title, subtitle, sections):
    doc = docx.Document()
    
    for sec in doc.sections:
        sec.top_margin = Inches(0.75)
        sec.bottom_margin = Inches(0.75)
        sec.left_margin = Inches(0.85)
        sec.right_margin = Inches(0.85)
        
    p_title = doc.add_paragraph()
    r_title = p_title.add_run(title)
    r_title.font.name = 'Calibri'
    r_title.font.size = Pt(18)
    r_title.font.bold = True
    r_title.font.color.rgb = RGBColor(30, 58, 138)
    
    p_sub = doc.add_paragraph()
    r_sub = p_sub.add_run(subtitle)
    r_sub.font.name = 'Calibri'
    r_sub.font.size = Pt(10.5)
    r_sub.font.italic = True
    r_sub.font.color.rgb = RGBColor(100, 116, 139)
    doc.add_paragraph().paragraph_format.space_after = Pt(6)
    
    for sec in sections:
        if sec.get('page_break_before', False):
            doc.add_page_break()
            
        h = doc.add_heading(sec['title'], level=2)
        h.paragraph_format.space_before = Pt(12)
        h.paragraph_format.space_after = Pt(3)
        for run in h.runs:
            run.font.name = 'Calibri'
            run.font.color.rgb = RGBColor(15, 23, 42)
            
        for p_text in sec.get('paragraphs', []):
            p = doc.add_paragraph(p_text)
            p.paragraph_format.space_after = Pt(5)
            for run in p.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(9.5)
                run.font.color.rgb = RGBColor(51, 65, 85)
                
    doc.save(filepath)


# -------------------------------------------------------------------------
# EML GENERATION HELPER
# -------------------------------------------------------------------------
def generate_branded_eml(filepath, from_email, to_email, subject, date_str, body_text):
    msg = EmailMessage()
    msg['From'] = from_email
    msg['To'] = to_email
    msg['Subject'] = subject
    msg['Date'] = date_str
    msg['Message-ID'] = f"<{datetime.now(timezone.utc).timestamp()}@mergeworks-dd-intake.com>"
    msg['MIME-Version'] = "1.0"
    msg.set_content(body_text)
    
    with open(filepath, 'wb') as f:
        f.write(msg.as_bytes())


# -------------------------------------------------------------------------
# SYNTHETIC MEDIA BUILDERS (MP3 & MP4)
# -------------------------------------------------------------------------
# These four files used to be downloaded from public sample URLs -- two
# Big Buck Bunny clips and an MDN heavy-metal track. That made the packets
# actively misleading: the ground truth described a founder interview, a CFO
# add-back call, a bottling plant tour and a Phase II site inspection, while
# the bytes on disk were cartoon footage and a metal song. A model could only
# "pass" those documents by hallucinating content that was not in the file, so
# the eval was rewarding exactly the failure mode it exists to catch.
#
# Media is now assembled from assets committed under test_sets/media_source:
#
#   scripts/   the interview / call / narration scripts -- source of truth,
#              written to satisfy each packet's existing ground truth
#   audio/     raw TTS renders of those scripts (multi-speaker for the two
#              dialogues, single narrator for the two walkthroughs)
#   stills/    facility stills used as the video frames
#
# Assembly is deterministic and fully offline: it needs only ffmpeg. It applies
# the acoustic character each file is supposed to have (conference-room tone for
# the in-person interview, narrowband phone-line for the recorded call, plant
# and site ambience under the walkthroughs) and builds the videos as moving-
# camera passes over the stills, which is what real site-visit footage is.
#
# Re-rendering the speech and stills from scratch needs the generation CLIs and
# is deliberately a separate opt-in step; see MEDIA_SOURCE_README.md.

MEDIA_SOURCE_DIR = os.path.join(ROOT_DIR, "test_sets", "media_source")

VIDEO_FPS = 25
VIDEO_W, VIDEO_H = 1280, 720


def _ffmpeg(args, label):
    """Run ffmpeg, surfacing the tail of stderr on failure."""
    proc = subprocess.run(["ffmpeg", "-y", "-v", "error"] + args,
                          capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"ffmpeg failed building {label}:\n{proc.stderr[-2000:]}")


def _media_duration(path):
    out = subprocess.run(
        ["ffprobe", "-v", "error", "-show_entries", "format=duration",
         "-of", "csv=p=0", path], capture_output=True, text=True).stdout
    return float(out.strip())


def _source_audio(name):
    path = os.path.join(MEDIA_SOURCE_DIR, "audio", name)
    if not os.path.exists(path):
        raise FileNotFoundError(
            f"Missing media source {path}. Speech assets ship with the repo; "
            "see test_sets/media_source/MEDIA_SOURCE_README.md to re-render.")
    return path


def build_room_audio(source_name, destination_path):
    """In-person recording: full band, slight room reverb, low HVAC floor."""
    src = _source_audio(source_name)
    _ffmpeg([
        "-i", src,
        "-f", "lavfi", "-i", "anoisesrc=color=brown:amplitude=0.004:r=44100",
        "-filter_complex",
        "[0:a]highpass=f=70,lowpass=f=13000,aecho=0.85:0.9:23:0.06,"
        "acompressor=threshold=-18dB:ratio=2.5:attack=20:release=250[v];"
        "[1:a]volume=0.4[n];"
        "[v][n]amix=inputs=2:duration=first:dropout_transition=0,"
        "alimiter=limit=0.95",
        "-map_metadata", "-1",
        "-ac", "1", "-ar", "44100", "-b:a", "128k", destination_path,
    ], os.path.basename(destination_path))
    print(f"Built {os.path.basename(destination_path)} "
          f"({_media_duration(destination_path):.1f}s room audio)")


def build_phone_call_audio(source_name, destination_path):
    """Recorded call: 300-3400 Hz narrowband, hard limiting, line noise."""
    src = _source_audio(source_name)
    _ffmpeg([
        "-i", src,
        "-f", "lavfi", "-i", "anoisesrc=color=pink:amplitude=0.008:r=16000",
        "-filter_complex",
        "[0:a]highpass=f=300,lowpass=f=3400,"
        "acompressor=threshold=-16dB:ratio=4:attack=5:release=120[v];"
        "[1:a]volume=0.55[n];"
        "[v][n]amix=inputs=2:duration=first:dropout_transition=0,"
        "alimiter=limit=0.95",
        "-map_metadata", "-1",
        "-ac", "1", "-ar", "16000", "-b:a", "48k", destination_path,
    ], os.path.basename(destination_path))
    print(f"Built {os.path.basename(destination_path)} "
          f"({_media_duration(destination_path):.1f}s narrowband call audio)")


def build_walkthrough_video(still_names, narration_name, destination_path,
                            ambience, handheld=False):
    """Assemble a narrated walkthrough: moving camera over stills + ambience.

    ambience: dict with noise `color`, `amp`, lowpass `lp` and mix `vol`.
    handheld: adds jitter, mild desaturation and sensor grain, for footage that
    is meant to read as a hand-carried inspection camera.
    """
    narration = _source_audio(narration_name)
    stills = [os.path.join(MEDIA_SOURCE_DIR, "stills", n) for n in still_names]
    missing = [s for s in stills if not os.path.exists(s)]
    if missing:
        raise FileNotFoundError(f"Missing stills: {missing}")

    # let the last shot breathe past the final word instead of cutting on it
    shot_seconds = (_media_duration(narration) + 1.6) / len(stills)
    frames = int(round(shot_seconds * VIDEO_FPS))

    with tempfile.TemporaryDirectory() as tmp:
        clips = []
        for i, still in enumerate(stills):
            clip = os.path.join(tmp, f"clip{i:02d}.mp4")
            # alternate a slow push-in and a slow pull-back
            zoom = (f"min(1.02+0.00050*on,1.14)" if i % 2 == 0
                    else f"max(1.14-0.00050*on,1.02)")
            vf = (
                "scale=1920:1080:force_original_aspect_ratio=increase,"
                "crop=1920:1080,"
                f"zoompan=z='{zoom}':d={frames}"
                f":x='iw/2-(iw/zoom/2)+{(-1) ** i}*70*sin(on/{frames}*3.14159)'"
                ":y='ih/2-(ih/zoom/2)'"
                f":s={VIDEO_W}x{VIDEO_H}:fps={VIDEO_FPS}"
            )
            if handheld:
                vf += (
                    f",scale={VIDEO_W + 48}:{VIDEO_H + 48},"
                    f"crop={VIDEO_W}:{VIDEO_H}"
                    ":'24+6*sin(n/9)+3*sin(n/23)':'24+5*sin(n/13)+3*cos(n/31)',"
                    "eq=saturation=0.92:contrast=1.03,noise=alls=5:allf=t"
                )
            # one input frame in, exactly `frames` frames out -- zoompan emits
            # d frames per input frame, so never hand it a looped stream
            _ffmpeg(["-i", still, "-vf", vf, "-frames:v", str(frames),
                     "-c:v", "libx264", "-preset", "veryfast", "-crf", "23",
                     "-pix_fmt", "yuv420p", "-r", str(VIDEO_FPS), clip],
                    os.path.basename(still))
            clips.append(clip)

        listfile = os.path.join(tmp, "concat.txt")
        with open(listfile, "w") as fh:
            for clip in clips:
                fh.write(f"file '{clip}'\n")
        silent = os.path.join(tmp, "silent.mp4")
        _ffmpeg(["-f", "concat", "-safe", "0", "-i", listfile, "-c", "copy",
                 silent], "concat")

        _ffmpeg([
            "-i", silent,
            "-i", narration,
            "-f", "lavfi", "-i",
            f"anoisesrc=color={ambience['color']}:"
            f"amplitude={ambience['amp']}:r=44100",
            "-filter_complex",
            "[1:a]highpass=f=90,lowpass=f=11000,"
            "acompressor=threshold=-18dB:ratio=2.5[v];"
            f"[2:a]lowpass=f={ambience['lp']},volume={ambience['vol']}[n];"
            "[v][n]amix=inputs=2:duration=first:dropout_transition=0,"
            "alimiter=limit=0.95[a]",
            "-map", "0:v", "-map", "[a]", "-map_metadata", "-1",
            "-c:v", "copy", "-c:a", "aac", "-b:a", "128k", "-ac", "1",
            "-movflags", "+faststart", "-shortest", destination_path,
        ], os.path.basename(destination_path))

    print(f"Built {os.path.basename(destination_path)} "
          f"({_media_duration(destination_path):.1f}s, {len(stills)} shots)")


PLANT_TOUR_STILLS = [
    "p4_shot1_dock.jpg", "p4_shot2_line1.jpg", "p4_shot3_line2_rebuild.jpg",
    "p4_shot4_line3.jpg", "p4_shot5_warehouse.jpg", "p4_shot6_wide_floor.jpg",
]
INSPECTION_STILLS = [
    "p6_shot1_tankfarm.jpg", "p6_shot2_cracked_berm.jpg",
    "p6_shot3_staining.jpg", "p6_shot4_monitoring_well.jpg",
    "p6_shot5_drums.jpg", "p6_shot6_wide_site.jpg",
]
# machine rumble under the plant floor; wind and open ground on the waste site
PLANT_AMBIENCE = {"color": "brown", "amp": "0.05", "lp": "900", "vol": "0.5"}
SITE_AMBIENCE = {"color": "pink", "amp": "0.04", "lp": "1400", "vol": "0.45"}


# =========================================================================
# PACKET 4: ATLANTIC BEVERAGE & BOTTLING CORP (PROCEED / GREEN)
# =========================================================================
def build_packet_4():
    p4_dir = os.path.join(DEALS_DIR, "packet_4_atlantic_beverage_proceed")
    os.makedirs(p4_dir, exist_ok=True)
    
    # 1. Multi-Tab Financial Model (8 Tabs) (.xlsx)
    wb = openpyxl.Workbook()
    
    # Tab 1: Executive KPI Dashboard
    ws_kpi = wb.active
    ws_kpi.title = "Executive KPI Dashboard"
    format_excel_tab(
        ws_kpi,
        title="Atlantic Beverage & Bottling Corp — Executive Summary & Operating KPIs",
        subtitle="Audited 5-Year Financial Track Record | Kentmere Ashford LLP",
        headers=["Performance Metric", "FY2021", "FY2022", "FY2023", "FY2024", "FY2025", "5-Yr CAGR", "Benchmark Target"],
        rows=[
            ["Net Case Delivery Volume (000s)", 1820, 1980, 2190, 2410, 2720, 0.105, ">2,500k Cases"],
            ["Gross Billed Revenue", 11800000, 12900000, 14200000, 16100000, 18500000, 0.119, ">10.0% Annual"],
            ["Gross Profit ($ USD)", 5664000, 6192000, 6816000, 7728000, 8880000, 0.119, "48.0% Margin"],
            ["Reported Operating EBITDA", 2242000, 2515000, 2803000, 3179000, 3650000, 0.130, ">18.0% Margin"],
            ["Normalized EBITDA Margin", 0.190, 0.195, 0.197, 0.197, 0.197, 0.009, "Industry Leading"],
            ["Free Cash Flow Conversion", 1680000, 1920000, 2180000, 2510000, 2940000, 0.150, ">80% Conversion"],
            ["Days Sales Outstanding (DSO)", 31.2, 30.8, 29.5, 28.4, 27.9, -0.027, "<35 Days"],
            ["Bottling Line Capacity Utilization", 0.48, 0.52, 0.58, 0.61, 0.65, 0.079, "35% Idle Capacity"]
        ],
        total_indices=[1, 2, 3, 5],
        is_currency_cols=[1, 2, 3, 4, 5]
    )
    
    # Tab 2: 5-Year Income Statement (40+ Lines)
    ws_pnl = wb.create_sheet("5-Year P&L Statement")
    format_excel_tab(
        ws_pnl,
        title="Atlantic Beverage & Bottling Corp — 5-Year Comprehensive Audited GAAP P&L",
        subtitle="Fiscal Years Ending Dec 31, 2021 through Dec 31, 2025 (in USD)",
        headers=["GAAP Account Category", "FY2021", "FY2022", "FY2023", "FY2024", "FY2025", "% Rev 2025"],
        rows=[
            ["Direct-Store-Delivery Bottled Beverages", 9440000, 10320000, 11360000, 12880000, 14800000, 0.800],
            ["Fountain Syrup & Co-Packing Contract Revenue", 2360000, 2580000, 2840000, 3220000, 3700000, 0.200],
            ["GROSS REVENUE", 11800000, 12900000, 14200000, 16100000, 18500000, 1.000],
            ["Raw Sugar & Concentrate Syrups", 3068000, 3354000, 3692000, 4186000, 4810000, 0.260],
            ["Aluminum Cans, Glass Bottles & Preforms", 1770000, 1935000, 2130000, 2415000, 2775000, 0.150],
            ["Direct Bottling Plant Machinist Labor", 708000, 774000, 852000, 966000, 1110000, 0.060],
            ["Freight-In & Raw Ingredient Packaging", 590000, 645000, 710000, 805000, 925000, 0.050],
            ["TOTAL COST OF GOODS SOLD", 6136000, 6708000, 7384000, 8372000, 9620000, 0.520],
            ["GROSS PROFIT", 5664000, 6192000, 6816000, 7728000, 8880000, 0.480],
            ["Route Driver Wages & Overtime", 1534000, 1677000, 1846000, 2093000, 2405000, 0.130],
            ["Refrigerated Fleet Fuel & Diesel", 354000, 387000, 426000, 483000, 555000, 0.030],
            ["Truck Fleet Maintenance & Insurance", 188800, 206400, 227200, 257600, 296000, 0.016],
            ["Warehouse Rent & Refrigeration Power", 590000, 645000, 710000, 805000, 925000, 0.050],
            ["Sales Representatives Commissions", 354000, 387000, 426000, 483000, 555000, 0.030],
            ["Corporate SG&A & Office Salaries", 472000, 516000, 568000, 644000, 740000, 0.040],
            ["TOTAL OPERATING EXPENSES", 3492800, 3818400, 4203200, 4765600, 5476000, 0.296],
            ["OPERATING INCOME (EBIT)", 2171200, 2373600, 2612800, 2962400, 3404000, 0.184],
            ["Bottling Line & Fleet Depreciation", 188800, 206400, 227200, 257600, 296000, 0.016],
            ["REPORTED EBITDA (AUDITED)", 2360000, 2580000, 2840000, 3220000, 3700000, 0.200],
            ["Normalized Owner Wage Replacement Delta", -118000, -65000, -37000, -41000, -50000, -0.003],
            ["TRUE AUDITED NORMALIZED EBITDA", 2242000, 2515000, 2803000, 3179000, 3650000, 0.197]
        ],
        total_indices=[2, 7, 8, 15, 16, 18, 20],
        is_currency_cols=[1, 2, 3, 4, 5]
    )
    
    # Tab 3: Balance Sheet (5 Years)
    ws_bs = wb.create_sheet("5-Year Balance Sheet")
    format_excel_tab(
        ws_bs,
        title="Atlantic Beverage & Bottling Corp — 5-Year Balance Sheet Track Record",
        subtitle="As of December 31, 2021 through 2025 (in USD)",
        headers=["Balance Sheet Line Item", "FY2021", "FY2022", "FY2023", "FY2024", "FY2025", "% Assets 2025"],
        rows=[
            ["Cash & Cash Equivalents", 820000, 940000, 1120000, 1280000, 1450000, 0.185],
            ["Accounts Receivable (Current <30 Days)", 980000, 1110000, 1260000, 1440000, 1620000, 0.207],
            ["Raw Beverage Syrup & Packaging Inventory", 520000, 610000, 710000, 810000, 890000, 0.114],
            ["TOTAL CURRENT ASSETS", 2320000, 2660000, 3090000, 3530000, 3960000, 0.506],
            ["High-Speed Automated Krones Bottling Line", 2100000, 2300000, 2500000, 2700000, 2850000, 0.364],
            ["Refrigerated DSD Delivery Fleet (22 Units)", 720000, 810000, 910000, 980000, 1020000, 0.130],
            ["TOTAL FIXED ASSETS (NET PP&E)", 2820000, 3110000, 3410000, 3680000, 3870000, 0.494],
            ["TOTAL ASSETS", 5140000, 5770000, 6500000, 7210000, 7830000, 1.000],
            ["Accounts Payable & Supplier Accruals", 710000, 820000, 930000, 1040000, 1120000, 0.143],
            ["Equipment Note (First National Bank)", 920000, 790000, 680000, 560000, 480000, 0.061],
            ["TOTAL LIABILITIES", 1630000, 1610000, 1610000, 1600000, 1600000, 0.204],
            ["Retained Earnings & Member Equity", 3510000, 4160000, 4890000, 5610000, 6230000, 0.796],
            ["TOTAL LIABILITIES & EQUITY", 5140000, 5770000, 6500000, 7210000, 7830000, 1.000]
        ],
        total_indices=[3, 6, 7, 10, 11, 12],
        is_currency_cols=[1, 2, 3, 4, 5]
    )
    
    # Tab 4: Customer Concentration & Master Account Roster (20+ Rows)
    ws_cust = wb.create_sheet("Customer Master & Roster")
    format_excel_tab(
        ws_cust,
        title="Atlantic Beverage — Customer Concentration & Contract Portfolio",
        subtitle="FY2025 Revenue Breakdown Across 142 Supermarket & Convenience Store Accounts",
        headers=["Customer Account Name", "Channel Type", "FY2025 Revenue", "% Share", "Contract Expiry", "Payment Terms", "Tenure"],
        rows=[
            ["ShopRite Supermarkets Regional Co-Op", "Supermarket Chain", 2127500, 0.115, "Dec 31, 2028", "Net 30", "14 Years"],
            ["Wawa Regional Convenience Stores (38)", "Convenience Store", 1850000, 0.100, "Oct 31, 2027", "Net 15", "11 Years"],
            ["Acme Markets Mid-Atlantic", "Supermarket Chain", 1665000, 0.090, "Jun 30, 2029", "Net 30", "9 Years"],
            ["QuickChek Convenience Group", "Convenience Store", 1295000, 0.070, "Mar 31, 2028", "Net 15", "8 Years"],
            ["Giant Food Stores Regional", "Supermarket Chain", 1110000, 0.060, "Dec 31, 2027", "Net 30", "6 Years"],
            ["Weis Markets Regional Network", "Supermarket Chain", 925000, 0.050, "Aug 31, 2028", "Net 30", "7 Years"],
            ["Independent Grocers Alliance (24 Outlets)", "Independent Retail", 1850000, 0.100, "Annual Evergreen", "Net 30", "12 Years"],
            ["Hospitality & University Dining (32 Accounts)", "Foodservice / Inst.", 2590000, 0.140, "Multi-Year Service", "Net 15", "10 Years"],
            ["Regional Beverage Distributors Co-Pack", "Contract Packing", 1850000, 0.100, "Dec 31, 2027", "Net 30", "5 Years"],
            ["Gas Station & Bodega DSD Network (50)", "Convenience Retail", 3232500, 0.175, "Weekly COD", "COD / Net 7", "15 Years"],
            ["TOTAL CONSOLIDATED REVENUE", "Consolidated", 18500000, 1.000, "Diversified Contract Base", "DSO = 27.9 Days", "Zero >11.5%"]
        ],
        total_indices=[10],
        is_currency_cols=[2]
    )
    
    # Tab 5: AR Aging & Bad Debt Ledger
    ws_ar = wb.create_sheet("AR Aging Ledger")
    format_excel_tab(
        ws_ar,
        title="Atlantic Beverage — Accounts Receivable Aging & Credit Quality",
        subtitle="As of December 31, 2025",
        headers=["Customer Category", "Current (0-30)", "31-60 Days", "61-90 Days", ">90 Days Past Due", "Total Outstanding", "Bad Debt Risk"],
        rows=[
            ["Major Supermarket Chains", 1120000, 45000, 0, 0, 1165000, "Low (Investment Grade)"],
            ["Regional Convenience Groups", 295000, 12000, 0, 0, 307000, "Low"],
            ["Independent Retailers", 115000, 18000, 4000, 1000, 138000, "Very Low (<0.7%)"],
            ["Foodservice & Hospitality", 90000, 10000, 0, 0, 100000, "Low"],
            ["TOTAL ACCOUNTS RECEIVABLE", 1620000, 85000, 4000, 1000, 1710000, "0.05% Historic Loss Rate"]
        ],
        total_indices=[4],
        is_currency_cols=[1, 2, 3, 4, 5]
    )
    
    xlsx_path = os.path.join(p4_dir, "Atlantic_Beverage_3Yr_Audited_Financial_Model.xlsx")
    wb.save(xlsx_path)
    
    # 2. Executive CIM (Multi-Page Styled PDF)
    pdf_path = os.path.join(p4_dir, "Atlantic_Beverage_Executive_CIM.pdf")
    generate_branded_pdf(
        pdf_path,
        title="Atlantic Beverage & Bottling Corp",
        subtitle="Confidential Information Memorandum (CIM) — Direct-Store-Delivery & Bottling Logistics",
        metadata_dict={
            "Enterprise Asking Price": "$18,500,000",
            "Normalized EBITDA (2025)": "$3,650,000 (19.7% Margin)",
            "5-Year Revenue CAGR": "11.9% ($11.8M to $18.5M)",
            "Imputed Transaction Multiple": "5.07x Normalized EBITDA",
            "Production Footprint": "45,000 Sq Ft Cold Warehouse",
            "Delivery Fleet Assets": "22 Refrigerated Freightliners"
        },
        sections=[
            {
                "heading": "1. Executive Summary & Investment Highlights",
                "paragraphs": [
                    "Atlantic Beverage & Bottling Corp is an established regional beverage bottler, contract packager, and Direct-Store-Delivery (DSD) distributor operating across the Mid-Atlantic corridor for over 15 years.",
                    "The business distributes high-velocity bottled beverages, craft functional sodas, and proprietary cold-brew teas to 142 premier retail grocery stores, regional convenience chains, and foodservice institutions.",
                    "With audited FY2025 revenue of $18,500,000 and Normalized EBITDA of $3,650,000, the company demonstrates pristine financial records audited annually by Kentmere Ashford LLP."
                ],
                "callout": {
                    "text": "<b>Investment Merits:</b><br/>• <b>Pristine Audit Quality:</b> 5 consecutive years of clean, unqualified GAAP audits.<br/>• <b>Low Customer Concentration:</b> No customer exceeds 11.5% of sales; top 5 accounts represent only 43.5%.<br/>• <b>Expansion Capacity:</b> Automated Krones high-speed bottling line operates at 65% capacity, enabling $10M+ in revenue growth without material CapEx.",
                    "type": "info"
                }
            },
            {
                "heading": "2. Financial Track Record & Historical Performance",
                "paragraphs": [
                    "Atlantic Beverage has generated strong top-line and EBITDA expansion driven by route density gains and direct-store-delivery economies of scale:"
                ],
                "table": {
                    "headers": ["Financial Metric", "FY2021", "FY2022", "FY2023", "FY2024", "FY2025", "CAGR"],
                    "rows": [
                        ["Gross Revenue", "$11,800,000", "$12,900,000", "$14,200,000", "$16,100,000", "$18,500,000", "+11.9%"],
                        ["Gross Profit", "$5,664,000", "$6,192,000", "$6,816,000", "$7,728,000", "$8,880,000", "+11.9%"],
                        ["Gross Margin", "48.0%", "48.0%", "48.0%", "48.0%", "48.0%", "Stable"],
                        ["Audited EBITDA", "$2,242,000", "$2,515,000", "$2,803,000", "$3,179,000", "$3,650,000", "+13.0%"],
                        ["EBITDA Margin", "19.0%", "19.5%", "19.7%", "19.7%", "19.7%", "+70 bps"]
                    ],
                    "colWidths": [120, 70, 70, 70, 70, 70, 70]
                }
            },
            {
                "heading": "3. Infrastructure & Distribution Logistics",
                "page_break_before": True,
                "paragraphs": [
                    "The business operates from a modern 45,000 square foot climate-controlled facility with 12 loading bays and dedicated racking for 2,400 pallets.",
                    "Manufacturing assets include an automated Krones 600-BPM bottling line (installed new in 2022) with automated CIP sanitation and packaging robotic cell.",
                    "The distribution fleet comprises 22 late-model Freightliner M2 refrigerated delivery trucks maintained under full-service OEM fleet contracts."
                ]
            }
        ]
    )
    
    # 3. Executed LOI Agreement (Multi-Page .docx)
    docx_path = os.path.join(p4_dir, "Atlantic_Beverage_Executed_LOI_Agreement.docx")
    generate_branded_docx(
        docx_path,
        title="BINDING LETTER OF INTENT FOR ASSET PURCHASE",
        subtitle="Transaction Agreement by and between Apex Buyout Partners and Atlantic Beverage & Bottling Corp",
        sections=[
            {
                "title": "SECTION 1. PURCHASE PRICE AND ENTERPRISE VALUATION",
                "paragraphs": [
                    "1.1 Total Purchase Price: The Buyer agrees to acquire substantially all operating assets of Atlantic Beverage & Bottling Corp for a total enterprise purchase price of $18,500,000 (Eighteen Million Five Hundred Thousand USD).",
                    "1.2 Valuation Basis: The purchase price is calculated as a 5.07x multiple applied to the agreed Normalized EBITDA of $3,650,000 for the trailing twelve-month period ending December 31, 2025.",
                    "1.3 Working Capital Adjustment: The purchase price includes a Target Net Working Capital peg of $2,840,000. At closing, the purchase price shall be adjusted dollar-for-dollar for any variance between the estimated closing NWC and the agreed target."
                ]
            },
            {
                "title": "SECTION 2. CONSIDERATION AT CLOSING & ESCROW TERMS",
                "paragraphs": [
                    "2.1 Cash Consideration: At closing, Buyer shall deliver $17,000,000 in immediate cash via federal wire transfer.",
                    "2.2 Seller Subordinated Note: Seller shall accept a $1,000,000 subordinated promissory note carrying 6.0% annual interest with quarterly amortization over a 36-month term.",
                    "2.3 Indemnity Escrow Deposit: An indemnity escrow of $500,000 (five hundred thousand USD) shall be funded from gross transaction proceeds and deposited with an independent escrow agent for a period of twelve (12) months to secure standard representations, warranties, and tax indemnifications."
                ]
            },
            {
                "title": "SECTION 3. TRANSITION SERVICES & NON-COMPETITION",
                "paragraphs": [
                    "3.1 Founder Transition Agreement: The founder and Chief Executive Officer shall execute a 12-month Transition Services Agreement (TSA) at an agreed market salary of $180,000 per annum to oversee route logistics continuity.",
                    "3.2 Non-Compete Covenant: Seller and key equity holders shall execute a 5-year, 150-mile geographic non-competition and non-solicitation agreement covering commercial beverage bottling and DSD distribution."
                ]
            },
            {
                "title": "SECTION 4. EXCLUSIVITY AND CONFIRMATORY DILIGENCE",
                "paragraphs": [
                    "4.1 Exclusivity: In consideration of Buyer incurring substantial legal, accounting, and environmental due diligence expenses, Seller agrees to grant Buyer sixty (60) days of strict exclusivity from the execution date hereof.",
                    "4.2 Access to Books and Records: Seller shall grant Buyer and its authorized representatives unrestricted access during normal business hours to all financial records, route delivery tracking data, customer contracts, and employee personnel files."
                ]
            }
        ]
    )
    
    # 4. Management Presentation (8 Slides .pptx)
    pptx_path = os.path.join(p4_dir, "Atlantic_Beverage_Management_Presentation.pptx")
    generate_branded_pptx(
        pptx_path,
        title="Atlantic Beverage & Bottling Corp",
        subtitle="Confidential Management Presentation — M&A Due Diligence Review",
        slides_data=[
            {
                "title": "Executive Summary: Premier Regional DSD Bottler",
                "bullets": [
                    "Founded in 2010; scaled to $18.5M revenue and $3.65M audited EBITDA across Mid-Atlantic markets.",
                    "Specialized in high-velocity Direct-Store-Delivery (DSD) to 142 supermarkets, convenience chains, and dining accounts.",
                    "5-year revenue CAGR of 11.9% with rock-solid 48.0% gross margin stability."
                ]
            },
            {
                "title": "Operational Infrastructure & Bottling Assets",
                "bullets": [
                    "45,000 sq ft modern climate-controlled warehouse featuring 12 loading docks and 2,400 pallet positions.",
                    "Automated Krones 600-BPM bottling and canning line installed in 2022 (OEM maintained).",
                    "Fleet of 22 Freightliner refrigerated trucks providing 99.4% on-time delivery across 14 weekly routes."
                ]
            },
            {
                "title": "Customer Diversification & Channel Breakdown",
                "bullets": [
                    "Zero single-customer concentration risk: Largest client (ShopRite) represents only 11.5% of revenue.",
                    "Top 5 accounts represent only 43.5% of total sales with average contract tenure exceeding 8.5 years.",
                    "DSD route density drives high barrier-to-entry against national competitors."
                ]
            },
            {
                "title": "Historical Financial & Margin Summary",
                "bullets": [
                    "Revenue expanded from $11.8M (2021) to $18.5M (2025) with zero customer churn among Tier-1 accounts.",
                    "Normalized EBITDA expanded to $3.65M (19.7% margin) with >80% free cash flow conversion.",
                    "Clean balance sheet with $1.45M cash, low debt (<0.5x leverage), and 27.9 DSO."
                ]
            },
            {
                "title": "Expansion Roadmap & Organic Growth Levers",
                "bullets": [
                    "Immediate opportunity to expand high-margin contract packing for functional beverages (+15% margin premium).",
                    "Krones bottling line operates at 65% capacity — ready to absorb $10M+ in additional volume without new CapEx.",
                    "Geographic bolt-on expansion into South Jersey and Northern Maryland route clusters."
                ]
            }
        ]
    )
    
    # 5. Broker Due Diligence Email Thread (.eml)
    eml_path = os.path.join(p4_dir, "Atlantic_Beverage_Broker_Due_Diligence_Thread.eml")
    generate_branded_eml(
        eml_path,
        from_email="diligence-lead@mergeworks.com",
        to_email="mark.stevens@atlanticcapitaladvisors.com",
        subject="Confirmatory Diligence Findings: Atlantic Beverage Maintenance CapEx & Key Employee Retention",
        date_str="Mon, 18 Aug 2026 14:22:10 -0400",
        body_text="""Hi Mark,

Our M&A due diligence team has concluded the on-site review and financial reconciliation for Atlantic Beverage & Bottling Corp:

1. Bottling Line Maintenance CapEx: The $296k annual D&A and maintenance reserves fully match the Krones OEM service history. No unrecorded deferred maintenance exists.
2. Route Driver & Manager Retention: All 22 route leads and the plant production supervisor have signed the 2-year retention bonus schedule.
3. Top Customer Renewals: ShopRite (11.5%) and Wawa (10.0%) contracts are verified through 2028 and 2027 respectively.

We confirm our full approval to proceed under the executed LOI terms ($18.5M enterprise valuation).

Best regards,
M&A Diligence Lead | MergeWorks"""
    )
    
    # 6. Founder / CFO management interview (.mp3, in-person room recording)
    audio_path = os.path.join(p4_dir, "Atlantic_Beverage_Founder_CFO_Interview.mp3")
    build_room_audio("p4_interview.mp3", audio_path)
    
    # 7. Bottling plant site walkthrough (.mp4, narrated)
    video_path = os.path.join(p4_dir, "Atlantic_Beverage_Bottling_Plant_Tour.mp4")
    build_walkthrough_video(PLANT_TOUR_STILLS, "p4_plant_tour_narration.mp3",
                            video_path, PLANT_AMBIENCE)


# =========================================================================
# PACKET 5: VANGUARD PRECISION AEROSPACE (RENEGOTIATE / YELLOW)
# =========================================================================
def build_packet_5():
    p5_dir = os.path.join(DEALS_DIR, "packet_5_vanguard_aerospace_renegotiate")
    os.makedirs(p5_dir, exist_ok=True)
    
    # 1. Multi-Tab Financial Model with Add-Back Ledger (6 Tabs) (.xlsx)
    wb = openpyxl.Workbook()
    
    # Tab 1: Valuation Bridge & Executive Summary
    ws_bridge = wb.active
    ws_bridge.title = "Valuation Bridge & Levers"
    format_excel_tab(
        ws_bridge,
        title="Vanguard Precision Aerospace — Forensic Valuation Waterfall & Re-Pricing",
        subtitle="Mandatory Purchase Price Renegotiation from $14.20M to $9.80M Target",
        headers=["Valuation Waterfall Step", "Reported Value", "Audit Disallowance", "Normalized Target", "Forensic Audit Rationale"],
        rows=[
            ["Seller Asking Price (CIM Claim)", 14200000, 0, 14200000, "Based on claimed EBITDA of $3.10M (4.58x Multiple)"],
            ["Disallow Personal Luxury Vehicle Leases", 0, -851000, -851000, "2x Porsche 911 leases on company payroll ($185k/yr)"],
            ["Disallow Family Personal Vacation Travel", 0, -667000, -667000, "Personal European trips expensed to SG&A ($145k/yr)"],
            ["Disallow Non-Working Son Salary", 0, -552000, -552000, "No operational role at company ($120k/yr)"],
            ["Pro-Forma Replacement GM Wage Reserve", 0, -920000, -920000, "CEO works 60 hrs/wk; replacement salary needed ($200k/yr)"],
            ["CNC Machinery Deferred CapEx Deduction", 0, -800000, -800000, "8x Haas 5-axis mills need immediate spindle overhauls"],
            ["Customer Concentration Risk Discount", 0, -610000, -610000, "38.5% single customer expiring in 4 mos under RFP"],
            ["TOTAL AUDIT PRICE REDUCTION", 0, -4400000, -4400000, "REQUIRED PRICE ADJUSTMENT: -$4,400,000 (-31.0%)"],
            ["FINAL TARGET ACQUISITION PRICE", 14200000, -4400000, 9800000, "4.56x True Audited Normalized EBITDA ($2,150,000)"]
        ],
        total_indices=[0, 7, 8],
        alert_indices=[1, 2, 3, 4, 5, 6, 7],
        is_currency_cols=[1, 2, 3]
    )
    
    # Tab 2: 5-Year Income Statement & Add-Back Normalization
    ws_pnl = wb.create_sheet("5-Year P&L & Add-Backs")
    format_excel_tab(
        ws_pnl,
        title="Vanguard Precision Aerospace — 5-Year P&L & Forensic Add-Back Analysis",
        subtitle="Fiscal Years Ending Dec 31, 2021 through 2025 (in USD)",
        headers=["Accounting Line Item", "FY2021", "FY2022", "FY2023", "FY2024", "FY2025", "% Rev 2025"],
        rows=[
            ["Precision Aerospace Machining Revenue", 11200000, 12100000, 13000000, 13800000, 14200000, 1.000],
            ["Titanium & Inconel Raw Alloy Billets", 4256000, 4598000, 4940000, 5244000, 5680000, 0.400],
            ["Direct CNC Machinist Wages & Overtime", 2688000, 2904000, 3120000, 3312000, 3408000, 0.240],
            ["GROSS PROFIT", 4256000, 4598000, 4940000, 5244000, 5112000, 0.360],
            ["CNC Tooling, Inserts & Coolant Wear", 560000, 610000, 680000, 740000, 780000, 0.055],
            ["Plant Lease & Heavy Power Utilities", 380000, 410000, 440000, 460000, 480000, 0.034],
            ["General Operating SG&A", 590000, 640000, 690000, 720000, 752000, 0.053],
            ["SELLER CLAIMED DISCRETIONARY ADD-BACKS", 310000, 360000, 410000, 430000, 450000, 0.032],
            ["SELLER CLAIMED ADJUSTED EBITDA", 2346000, 2594000, 2800000, 2982000, 3100000, 0.218],
            ["Disallow Personal Luxury Vehicle Leases", -120000, -140000, -160000, -175000, -185000, -0.013],
            ["Disallow Family Personal Travel", -100000, -110000, -130000, -140000, -145000, -0.010],
            ["Disallow Non-Working Son Salary", -90000, -110000, -120000, -115000, -120000, -0.008],
            ["Replacement GM Management Salary", -160000, -170000, -180000, -190000, -200000, -0.014],
            ["Unfunded CNC Spindle Repair Reserve", -180000, -210000, -250000, -280000, -300000, -0.021],
            ["TRUE AUDITED NORMALIZED EBITDA", 1696000, 1854000, 1960000, 2082000, 2150000, 0.151]
        ],
        total_indices=[3, 8, 14],
        alert_indices=[7, 9, 10, 11, 12, 13, 14],
        is_currency_cols=[1, 2, 3, 4, 5]
    )
    
    # Tab 3: Detailed Add-Back Ledger (15 Line Items)
    ws_add = wb.create_sheet("Itemized Add-Back Ledger")
    format_excel_tab(
        ws_add,
        title="Vanguard Precision Aerospace — Itemized Seller Add-Back Audit Ledger",
        subtitle="Forensic Vouching Results for FY2025",
        headers=["Add-Back Category", "Seller Claim", "Audited Legitimacy", "Disallowed Amount", "Auditor Vouching Comments"],
        rows=[
            ["Porsche 911 Turbo Leases (2x)", 185000, "DISALLOWED (0%)", 185000, "Personal luxury vehicles used by owner's spouse and child"],
            ["European Summer Vacation Flights", 145000, "DISALLOWED (0%)", 145000, "Personal family vacation charged to company corporate Amex"],
            ["Consulting Salary to Son (Jason)", 120000, "DISALLOWED (0%)", 120000, "Enrolled full-time at university; zero work product produced"],
            ["Country Club Golf Memberships", 38000, "DISALLOWED (0%)", 38000, "Non-business recreational membership"],
            ["Personal Home Generator Installation", 24000, "DISALLOWED (0%)", 24000, "Expensed as plant utility repair; installed at owner's residence"],
            ["One-Time AS9100D Recertification", 45000, "APPROVED (100%)", 0, "Legitimate non-recurring audit preparation expense"],
            ["Non-Recurring Scrap Part Settlement", 35000, "APPROVED (100%)", 0, "Isolated tooling defect settlement with supplier"],
            ["TOTAL ADD-BACK SCHEDULE", 592000, "DISALLOWED: $512k", 512000, "OVER 86% OF SELLER ADD-BACK CLAIMS ARE FABRICATED"]
        ],
        total_indices=[7],
        alert_indices=[0, 1, 2, 3, 4, 7],
        is_currency_cols=[1, 3]
    )
    
    # Tab 4: Machinery Fleet Inspection & CapEx Backlog (24 Machines)
    ws_mach = wb.create_sheet("CNC Machinery & CapEx Backlog")
    format_excel_tab(
        ws_mach,
        title="Vanguard Precision Aerospace — CNC Machine Fleet Health & Overhaul Costs",
        subtitle="Independent Assessment by Precision Tooling Diagnostics LLC",
        headers=["Machine Tag", "Model & Spec", "Year", "Spindle Runout", "Guideway Condition", "Overhaul Cost ($)", "Status"],
        rows=[
            ["CNC-01 to CNC-04", "Haas VF-4 5-Axis (4 Units)", 2012, "0.0011 in (FAIL)", "Heavy Backlash", 160000, "Immediate Rebuild Required"],
            ["CNC-05 to CNC-08", "Haas VF-6 5-Axis (4 Units)", 2014, "0.0009 in (FAIL)", "Moderate Wear", 160000, "Immediate Rebuild Required"],
            ["CNC-09 to CNC-10", "Makino A61 Horizontal (2)", 2015, "0.0007 in (FAIL)", "Hydraulic Leak", 160000, "Hydraulic Unit Replacement"],
            ["CNC-11 to CNC-16", "Mori Seiki NV5000 (6)", 2018, "0.0004 in (PASS)", "Good", 80000, "Preventative Service"],
            ["CNC-17 to CNC-24", "Doosan Puma Lathes (8)", 2020, "0.0002 in (PASS)", "Excellent", 40000, "Standard Maintenance"],
            ["TOTAL MACHINERY FLEET", "24 Active CNC Centers", 2016, "8 of 24 Out of Spec", "High Deferred Wear", 600000, "$600k-$800k CapEx Backlog"]
        ],
        total_indices=[5],
        alert_indices=[0, 1, 2, 5],
        is_currency_cols=[5]
    )
    
    xlsx_path = os.path.join(p5_dir, "Vanguard_Aerospace_Monthly_PnL_AddBack_Ledger.xlsx")
    wb.save(xlsx_path)
    
    # 2. Customer Concentration AR Master (.xlsx)
    wb_ar = openpyxl.Workbook()
    ws_ar = wb_ar.active
    ws_ar.title = "Customer Concentration Master"
    format_excel_tab(
        ws_ar,
        title="Vanguard Precision Aerospace — Customer Concentration & Contract Master",
        subtitle="FY2025 Revenue Breakdown Across Defense & Aerospace Programs",
        headers=["Customer Account", "Program Platform", "FY2025 Revenue", "% Share", "Contract Expiration", "Re-Compete Risk"],
        rows=[
            ["Fabrikam Aerostructures", "AX-114 Structural Brackets", 5467000, 0.385, "April 30, 2026", "CRITICAL: 15% Rate Cut RFP Notice Issued"],
            ["Adatum Guidance Systems", "Actuator Guidance Housings", 2556000, 0.180, "Dec 31, 2027", "Moderate: Active delivery schedule"],
            ["Contoso Land Systems", "Titanium Fairings", 1846000, 0.130, "Nov 30, 2026", "Low: Multi-year blanket PO"],
            ["Proseware Precision Tooling", "Precision Fixtures", 1420000, 0.100, "Annual Evergreen", "Low: Long-standing vendor status"],
            ["Tailspin Defense Group", "Rotor Hub Flanges", 1136000, 0.080, "Oct 31, 2026", "Moderate"],
            ["Commercial Machine Accounts (12)", "Commercial Tooling", 1775000, 0.125, "Job Shop POs", "Diversified"],
            ["TOTAL CONSOLIDATED REVENUE", "Consolidated", 14200000, 1.000, "Top Customer = 38.5%", "HIGH CONCENTRATION VULNERABILITY"]
        ],
        total_indices=[6],
        alert_indices=[0, 6],
        is_currency_cols=[2]
    )
    ar_path = os.path.join(p5_dir, "Vanguard_Aerospace_Customer_Concentration_AR.xlsx")
    wb_ar.save(ar_path)
    
    # 3. Confidential Information Memorandum (Multi-Page PDF)
    pdf_path = os.path.join(p5_dir, "Vanguard_Aerospace_CIM_Executive_Overview.pdf")
    generate_branded_pdf(
        pdf_path,
        title="Vanguard Precision Aerospace Machining",
        subtitle="Confidential Executive Overview & Acquisition Prospectus",
        metadata_dict={
            "Seller Asking Price": "$14,200,000 (DEMAND)",
            "Seller Claimed EBITDA": "$3,100,000 (UNRELIABLE)",
            "Audited Normalized EBITDA": "$2,150,000 (-30.6% Deficit)",
            "Mandatory Target Valuation": "$9,800,000 (-$4.40M Cut)",
            "Top Customer Concentration": "38.5% (Fabrikam Aerostructures)",
            "Deferred Machinery Backlog": "$800,000 Spindle Repairs"
        },
        sections=[
            {
                "heading": "1. Valuation Re-Pricing Directive (RENEGOTIATE)",
                "paragraphs": [
                    "Vanguard Precision Aerospace Machining operates 24 CNC machining centers producing high-precision titanium and structural aluminum components for major defense prime contractors.",
                    "While the company maintains active AS9100D and ITAR certifications, forensic audit reveals three severe valuation impairments that mandate an immediate purchase price reduction from $14.2M to $9.8M."
                ],
                "callout": {
                    "text": "<b>RE-PRICING WATERFALL LEVERS:</b><br/>• <b>Add-Back Disallowances (-$2.07M):</b> $450k/yr of personal luxury vehicles, family vacations, and unearned family salaries disallowed.<br/>• <b>Customer Concentration Risk (-$610k):</b> Top customer (38.5%) contract expires in 4 months under a mandatory 15% rate cut RFP.<br/>• <b>Machinery CapEx Backlog (-$800k):</b> 8x Haas 5-axis mills require immediate spindle and guideway rebuilds.",
                    "type": "warning"
                }
            },
            {
                "heading": "2. Forensic Valuation Waterfall",
                "table": {
                    "headers": ["Bridge Component", "Valuation Impact", "Multiple", "Resulting Price"],
                    "rows": [
                        ["Seller Asking Price", "$14,200,000", "4.58x Claimed", "$14,200,000"],
                        ["Disallow $450k Personal Add-Backs", "-$2,070,000", "4.60x Multiple", "$12,130,000"],
                        ["Replacement GM & Wage Adjustments", "-$920,000", "4.60x Multiple", "$11,210,000"],
                        ["CapEx Deferred Maintenance Deduction", "-$800,000", "Direct Offset", "$10,410,000"],
                        ["Customer Concentration Discount (38.5%)", "-$610,000", "Risk Discount", "$9,800,000"],
                        ["FINAL TARGET TRANSACTION PRICE", "-$4,400,000", "4.56x Audited", "$9,800,000"]
                    ],
                    "colWidths": [160, 110, 110, 150]
                }
            }
        ]
    )
    
    # 4. Engineering Machinery Fleet Inspection (Multi-Page .docx)
    docx_path = os.path.join(p5_dir, "Vanguard_Aerospace_Machinery_Fleet_Inspection.docx")
    generate_branded_docx(
        docx_path,
        title="TECHNICAL MACHINERY INSPECTION & FLEET DIAGNOSTIC REPORT",
        subtitle="Independent Assessment of 24 CNC Machining Centers | Precision Tooling Diagnostics LLC",
        sections=[
            {
                "title": "1. SCOPE AND METHODOLOGY OF INSPECTION",
                "paragraphs": [
                    "Precision Tooling Diagnostics LLC performed an exhaustive on-site mechanical, electrical, and laser interferometry diagnostic inspection of the 24 CNC machining centers at Vanguard Precision Aerospace on January 14-16, 2026.",
                    "Testing included Renishaw ballbar dynamic circularity testing, spindle dynamic runout measurement, hydraulic power unit fluid analysis, and linear scale backlash calibration."
                ]
            },
            {
                "title": "2. CRITICAL MACHINERY FAILURES & DEFECTS",
                "paragraphs": [
                    "Eight (8) Haas 5-axis vertical machining centers (CNC-01 through CNC-08, vintage 2012-2014) exhibited severe spindle bearing degradation with total indicated runout (TIR) measuring 0.0009 in to 0.0011 in, exceeding AS9100D aerospace tolerance limits (max 0.0002 in).",
                    "Two (2) Makino A61 horizontal machining centers require full hydraulic power unit rebuilds and way-cover replacements due to chronic high-pressure coolant ingress.",
                    "Estimated total cost for replacement spindles, ball screws, and optical scale calibration across the fleet is $600,000 to $800,000."
                ]
            },
            {
                "title": "3. BUYER TRANSACTION RECOMMENDATION",
                "paragraphs": [
                    "Buyer must demand a direct $800,000 reduction from the cash purchase price or require Seller to escrow $800,000 to complete all spindle overhauls prior to the closing date."
                ]
            }
        ]
    )
    
    # 5. Top Customer Contract Notice Email (.eml)
    eml_path = os.path.join(p5_dir, "Vanguard_Aerospace_Top_Customer_Contract_Notice.eml")
    generate_branded_eml(
        eml_path,
        from_email="procurement.aero@fabrikam-aero.com",
        to_email="j.vanguard@vanguardaero.com",
        subject="URGENT: Long-Term Agreement Expiration Notice & Mandatory 15% RFP Re-Compete",
        date_str="Wed, 07 Jan 2026 11:15:40 -0500",
        body_text="""Dear Mr. Vanguard,

This letter serves as formal notification that Long-Term Pricing Agreement #FBK-AERO-4491 for AX-114 structural titanium mounting brackets expires on April 30, 2026.

As part of Fabrikam's enterprise cost-reduction initiative, this component package is being released to competitive RFP bidding among four certified suppliers. To remain in consideration as primary supplier, Vanguard must submit a revised pricing schedule reflecting an across-the-board 15% unit cost reduction.

Failure to meet competitive bidding benchmarks may result in dual-sourcing or full reallocation of volume starting Q3 2026.

Sincerely,
Global Aerospace Procurement | Fabrikam Aerostructures Corporation"""
    )
    
    # 6. CFO add-back clarification call (.mp3, recorded phone line).
    #    Carries this packet's modality-exclusive yellow flag: the CFO
    #    confirms personal luxury vehicle leases inside the seller add-backs.
    #    That admission appears in no document in the packet.
    audio_path = os.path.join(p5_dir, "Vanguard_Aerospace_CFO_Clarification_Call.mp3")
    build_phone_call_audio("p5_call.mp3", audio_path)


# =========================================================================
# PACKET 6: TERRACLEAN INDUSTRIAL WASTE (WALK AWAY / RED)
# =========================================================================
def build_packet_6():
    p6_dir = os.path.join(DEALS_DIR, "packet_6_terraclean_waste_walkaway")
    os.makedirs(p6_dir, exist_ok=True)
    
    # 1. Seller CIM Overview (Multi-Page PDF)
    pdf_path = os.path.join(p6_dir, "TerraClean_Waste_Seller_CIM_Overview.pdf")
    generate_branded_pdf(
        pdf_path,
        title="TerraClean Industrial Waste Solutions LLC",
        subtitle="Confidential Acquisition Teaser & Forensic Due Diligence Assessment",
        metadata_dict={
            "Seller Asking Price": "$22,000,000",
            "Claimed Revenue": "$22,000,000 (FABRICATED)",
            "Verified Bank Cash": "$11,200,000 ($10.8M Deficit)",
            "True Operating Income": "-$350,000 (INSOLVENT)",
            "Undisclosed Liabilities": "$4,200,000 EPA Superfund Lien",
            "Delinquent AR (>180 Days)": "$4,800,000 (71.9% Bad Debt)"
        },
        sections=[
            {
                "heading": "1. FATAL DEAL-BREAKER: Immediate Walk-Away Recommendation",
                "paragraphs": [
                    "TerraClean Industrial Waste Solutions advertises itself as a leading hazardous waste treatment, contaminated soil remediation, and industrial chemical disposal contractor.",
                    "Forensic audit across banking records, IRS Form 1120 tax filings, EPA regulatory dockets, and employee whistleblower disclosures uncovers catastrophic structural fraud, phantom round-trip revenue transfers, and severe environmental contamination liabilities."
                ],
                "callout": {
                    "text": "<b>FATAL DEAL-BREAKER (WALK AWAY IMMEDIATELY):</b><br/>• <b>Phantom Revenue ($10.8M Deficit):</b> Seller claims $22.0M revenue vs verified bank revenue of only $11.20M (49.1% inflation).<br/>• <b>EPA Superfund Consent Decree:</b> Active $4.2M federal CERCLA cleanup lien and 60-day operating permit cure warning.<br/>• <b>74% Uncollectible AR:</b> $4.8M of accounts receivable is >180 days delinquent with bankrupt contractors.",
                    "type": "danger"
                }
            },
            {
                "heading": "2. Forensic Discrepancy Matrix",
                "table": {
                    "headers": ["Financial Metric", "CIM Claimed", "Verified Tax / Bank", "Variance", "Forensic Severity"],
                    "rows": [
                        ["Gross Revenue", "$22,000,000", "$11,200,000", "-$10,800,000", "CRITICAL FRAUD (49.1%)"],
                        ["Gross Profit", "$12,100,000", "$3,360,000", "-$8,740,000", "True Margin 30% vs 55%"],
                        ["Operating EBITDA", "$4,500,000", "-$350,000", "-$4,850,000", "Operational Insolvency"],
                        ["Aged AR (>180 Days)", "$0 Reported", "$4,800,000", "+$4,800,000", "71.9% Bad Debt Loss"],
                        ["Environmental Liabilities", "Zero Disclosed", "$4,200,000", "+$4,200,000", "Federal CERCLA Lien"]
                    ],
                    "colWidths": [130, 95, 105, 105, 105]
                }
            }
        ]
    )
    
    # 2. General Ledger vs Bank Reconciliation Discrepancy (4 Tabs) (.xlsx)
    wb_bank = openpyxl.Workbook()
    ws_bank = wb_bank.active
    ws_bank.title = "GL vs Bank Reconciliation"
    format_excel_tab(
        ws_bank,
        title="TerraClean Waste Solutions — Bank Statement vs GL Reconciliation Discrepancy",
        subtitle="FY2025 Monthly Bank Deposits vs Reported General Ledger Invoicing",
        headers=["Month / Period", "Reported GL Invoicing", "Verified Bank Deposits", "Phantom Round-Trip Deposits", "Net Real Cash Discrepancy"],
        rows=[
            ["January 2025", 1750000, 920000, -450000, -830000],
            ["February 2025", 1800000, 950000, -500000, -850000],
            ["March 2025", 1650000, 880000, -450000, -770000],
            ["April 2025", 1900000, 980000, -520000, -920000],
            ["May 2025", 1850000, 960000, -480000, -890000],
            ["June 2025", 1850000, 960000, -500000, -890000],
            ["July 2025", 1800000, 940000, -480000, -860000],
            ["August 2025", 1750000, 910000, -460000, -840000],
            ["September 2025", 1850000, 950000, -510000, -900000],
            ["October 2025", 1950000, 930000, -490000, -1020000],
            ["November 2025", 1900000, 920000, -480000, -980000],
            ["December 2025", 1950000, 900000, -480000, -1050000],
            ["TOTAL FY2025", 22000000, 11200000, -5800000, -10800000]
        ],
        total_indices=[12],
        alert_indices=list(range(13)),
        is_currency_cols=[1, 2, 3, 4]
    )
    bank_path = os.path.join(p6_dir, "TerraClean_Waste_General_Ledger_Bank_Mismatch.xlsx")
    wb_bank.save(bank_path)
    
    # 3. IRS Form 1120 Federal Tax Return Audit (3 Tabs) (.xlsx)
    wb_tax = openpyxl.Workbook()
    ws_tax = wb_tax.active
    ws_tax.title = "Form 1120 Tax Audit"
    format_excel_tab(
        ws_tax,
        title="TerraClean Waste Solutions — IRS Form 1120 Corporate Tax Audit",
        subtitle="Filed with the Internal Revenue Service | Fiscal Year Ending Dec 31, 2025",
        headers=["IRS Form 1120 Line Item", "Filed Tax Amount ($)", "CIM Claimed Amount ($)", "Variance / Discrepancy", "Audit Finding"],
        rows=[
            ["Line 1a: Gross Receipts or Sales", 11200000, 22000000, -10800000, "FATAL: 49.1% Revenue Inflation"],
            ["Line 2: Cost of Goods Sold & Landfill Tipping", 7840000, 9900000, -2060000, "Severe disposal fee underreporting"],
            ["Line 3: Gross Profit", 3360000, 12100000, -8740000, "True Gross Margin: 30.0% vs 55.0%"],
            ["Line 12: Compensation of Officers & Wages", 2150000, 3200000, -1050000, "Hazardous technician payroll"],
            ["Line 26: Other Deductions & Legal Defense", 1560000, 4400000, -2840000, "Excludes environmental fines"],
            ["Line 28: Total Deductions", 3710000, 7600000, -3890000, "Overhead expenses"],
            ["Line 30: TAXABLE OPERATING INCOME (LOSS)", -350000, 4500000, -4850000, "TRUE OPERATING INSOLVENCY (-$350k)"]
        ],
        total_indices=[2, 5, 6],
        alert_indices=[0, 2, 6],
        is_currency_cols=[1, 2, 3]
    )
    tax_path = os.path.join(p6_dir, "TerraClean_Waste_Form_1120_Federal_Tax_Audit.xlsx")
    wb_tax.save(tax_path)
    
    # 4. AR Aging & Bad Debt Ledger (3 Tabs) (.xlsx)
    wb_ar = openpyxl.Workbook()
    ws_ar = wb_ar.active
    ws_ar.title = "AR Aging & Bad Debt"
    format_excel_tab(
        ws_ar,
        title="TerraClean Waste Solutions — Accounts Receivable Delinquency & Bad Debt Ledger",
        subtitle="As of December 31, 2025",
        headers=["Customer Account / Debtor", "Current (0-60)", "61-120 Days", ">180 Days Delinquent", "Collectibility Assessment"],
        rows=[
            ["Apex Chemical Remediation Corp", 0, 0, 1850000, "UNCOLLECTIBLE: Chapter 7 Liquidation Filed"],
            ["Tri-State Environmental Contractors", 0, 120000, 1420000, "UNCOLLECTIBLE: License Revoked, Disputed Billing"],
            ["Midwest Soil Recovery LLC", 0, 80000, 1530000, "UNCOLLECTIBLE: Unresponsive to Legal Demand (>240 Days)"],
            ["Active Municipal Treatment Accounts", 1420000, 260000, 0, "Collectibility Normal"],
            ["TOTAL ACCOUNTS RECEIVABLE", 1420000, 460000, 4800000, "TOTAL AR: $6.68M ($4.80M Bad Debt = 71.9% Loss)"]
        ],
        total_indices=[4],
        alert_indices=[0, 1, 2, 4],
        is_currency_cols=[1, 2, 3]
    )
    ar_path = os.path.join(p6_dir, "TerraClean_Waste_AR_Aging_Bad_Debt_Ledger.xlsx")
    wb_ar.save(ar_path)
    
    # 5. EPA Superfund Consent Decree (Multi-Page .docx)
    docx_path = os.path.join(p6_dir, "TerraClean_Waste_EPA_Superfund_Consent_Decree.docx")
    generate_branded_docx(
        docx_path,
        title="CONFIDENTIAL LEGAL DISCLOSURE: EPA SUPERFUND CONSENT DECREE",
        subtitle="United States District Court — Civil Action No. 25-CV-8891-EPA",
        sections=[
            {
                "title": "SECTION 1. NATURE OF FEDERAL ENFORCEMENT ACTION",
                "paragraphs": [
                    "The United States Environmental Protection Agency (EPA), in conjunction with the United States Department of Justice (DOJ), entered a formal Consent Decree on October 14, 2025, against TerraClean Industrial Waste Solutions LLC under Section 106 and 107 of the Comprehensive Environmental Response, Compensation, and Liability Act (CERCLA), 42 U.S.C. §§ 9606, 9607.",
                    "TerraClean is subject to an active $4,200,000 joint and several remediation obligation for soil and groundwater chlorinated solvent contamination at the Newark Treatment Facility."
                ]
            },
            {
                "title": "SECTION 2. STATE OPERATING PERMIT REVOCATION WARNING",
                "paragraphs": [
                    "The State Department of Environmental Protection issued a 60-day Notice of Intent to Revoke Part B Hazardous Waste Storage and Treatment Permit #HW-00921 on December 1, 2025, citing repeated secondary containment lining failures.",
                    "Failure to cure will result in immediate facility closure and revocation of operational licenses by March 31, 2026."
                ]
            },
            {
                "title": "SECTION 3. SUCCESSOR LIABILITY AND BUYER WARNING",
                "paragraphs": [
                    "Under federal CERCLA statutes, any purchaser of company assets will assume strict, joint, and several successor liability for all environmental contamination. The transaction presents insurmountable insolvency risk. WALK AWAY IMMEDIATELY."
                ]
            }
        ]
    )
    
    # 6. Whistleblower Internal Email (.eml)
    eml_path = os.path.join(p6_dir, "TerraClean_Waste_Whistleblower_Internal_Email.eml")
    generate_branded_eml(
        eml_path,
        from_email="lab-director@terracleanwaste.com",
        to_email="exec-team@terracleanwaste.com",
        subject="URGENT: Falsified Waste Manifests and Secondary Containment Leaks (Tank Farm B)",
        date_str="Fri, 19 Dec 2025 18:42:01 -0500",
        body_text="""Management,

I am writing to create a formal written record of our laboratory testing results.

The chlorinated solvent manifest logs submitted to the state inspector last week did NOT reflect the true toxicity levels in holding Tank 4B. The secondary containment lining has failed, and solvent leachate is seeping into the perimeter monitoring well.

If the prospective buyer's environmental Phase II engineers perform soil core drilling anywhere near the eastern retention basin, they will immediately discover chlorinated solvent contamination exceeding federal limits by 800x.

We cannot continue billing customers for certified destruction when the waste is merely being blended and stockpiled.

Compliance Director | TerraClean Waste"""
    )
    
    # 7. Phase II site inspection footage (.mp4, handheld inspection camera).
    #    Carries this packet's modality-exclusive red flag: structural
    #    failure of secondary containment at Tank Farm B.
    video_path = os.path.join(p6_dir, "TerraClean_Waste_Facility_Inspection_Site.mp4")
    build_walkthrough_video(INSPECTION_STILLS, "p6_inspection_narration.mp3",
                            video_path, SITE_AMBIENCE, handheld=True)


if __name__ == "__main__":
    print("Building Deep Multi-Page Packet 4: Atlantic Beverage & Bottling (PROCEED / GREEN)...")
    build_packet_4()
    print("Building Deep Multi-Page Packet 5: Vanguard Precision Aerospace (RENEGOTIATE / YELLOW)...")
    build_packet_5()
    print("Building Deep Multi-Page Packet 6: TerraClean Industrial Waste (WALK AWAY / RED)...")
    build_packet_6()
    print("ALL 3 DEEP MULTI-PAGE DOSSIERS & REAL MEDIA PACKETS SUCCESSFULLY GENERATED!")
